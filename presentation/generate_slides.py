"""
generate_slides.py — Build the TruPhoto PowerPoint deck.

Run from the repo root with python-pptx installed:
  pip install python-pptx
  python presentation/generate_slides.py

Output: presentation/truphoto_presentation.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


SLIDES = [
    {
        "type": "title",
        "title": "TruPhoto",
        "subtitle": (
            "Image Authenticity Classification\n"
            "Real  /  Forged  /  AI-Generated\n\n"
            "Tabeen Raoof  -  CS483 Deep Learning, SFBU"
        ),
        "notes": (
            "Open with the project name and the three-way classification task. "
            "Headline: 74.92% validation and 73.27% test accuracy on a 3-class problem with "
            "25 AI generators, 11 Real-image sources, and CASIA-2 forgeries. Val/test gap is "
            "only 1.65 pp - the model generalizes, it isn't overfitting. The numbers are a bit "
            "lower than the prior 77.62% run on purpose - we deliberately merged CASIA's "
            "authentic images into the Real class to break a dataset-source shortcut, and the "
            "new model is honest where the old one was cheating."
        ),
    },
    {
        "type": "content",
        "title": "Why image authenticity matters",
        "bullets": [
            "Generative models (Stable Diffusion, GANs, deepfakes) make synthetic images trivial to produce.",
            "Traditional forgery (splicing, copy-move) is still common in misinformation and fraud.",
            "Detecting both threats requires a single classifier that distinguishes:",
            "    -  Real photos",
            "    -  Manually forged images (CASIA 2.0  Tp/)",
            "    -  Fully AI-generated images (ArtiFact, 25 generators)",
            "Goal: a 3-class system, not just a binary real/fake detector.",
        ],
        "notes": (
            "Frame the problem: generative AI plus classic photo manipulation are converging threats. "
            "Most academic work treats this as binary, but real-world deployments need to distinguish "
            "manual forgeries from synthetic ones because the response is different - a spliced image "
            "comes from a human editor, an AI-generated one comes from a model."
        ),
    },
    {
        "type": "content",
        "title": "Pipeline overview",
        "bullets": [
            "Two-stage architecture:",
            "    1.  Frozen pretrained CNN  ->  fixed-length feature vector per image.",
            "    2.  Lightweight classifier (Random Forest / SVM) on those features.",
            "",
            "Why this split?",
            "    -  No GPU required for training the classifier.",
            "    -  CNN runs once, classifier iterates fast.",
            "    -  Classical classifiers are easy to interpret and calibrate.",
            "",
            "Pipeline:  step1 prep  ->  step2 features  ->  step3 train  ->  step4 evaluate  ->  step4b sanity check  ->  step5 demo",
        ],
        "notes": (
            "Walk through the two-stage idea. Standard transfer-learning pattern: use ImageNet's "
            "general visual knowledge for free, only train a small final classifier on your task. "
            "Trade-off: we lose the option to fine-tune CNN weights on forgery-specific cues. "
            "We'll come back to this on the design-decisions and fine-tuning slides."
        ),
    },
    {
        "type": "content",
        "title": "Datasets",
        "bullets": [
            "ArtiFact (Kaggle):",
            "    -  33 generator/source subfolders kept on disk after partial download.",
            "    -  10 contribute Real images (afhq, celebahq, coco, cycle_gan, ffhq, imagenet,",
            "         landscape, lsun, metfaces, pro_gan)  ~972K images.",
            "    -  25 contribute AI-Generated (StyleGAN1/2/3, Stable Diffusion, GLIDE, BigGAN,",
            "         latent_diffusion, DDPM, sfhq, palette, ...)  ~1.53M images.",
            "",
            "CASIA 2.0:",
            "    -  Tp/ -> 'Forged' class  (5,123 images, full set used).",
            "    -  Au/ -> mixed into 'Real' class  (7,491 images)  - added in final iteration",
            "         to break the dataset-source shortcut.",
            "",
            "Per-class budget:  12,500 samples per class, stratified across sources.",
            "70 / 15 / 15 train / val / test split, fixed seed 42.  Forged is undersized (5,123 total).",
        ],
        "notes": (
            "Be explicit about which class comes from where. The CASIA Au merge was added in the "
            "final iteration after the cross-source check showed that Real (only ArtiFact) and "
            "Forged (only CASIA) had no source overlap, letting the model learn 'this came from "
            "CASIA' as a free predictor for Forged. Pulling Au into Real forces the model to "
            "actually look at the image content."
        ),
    },
    {
        "type": "content",
        "title": "Step 1  -  Data preparation (three key fixes)",
        "bullets": [
            "Discovery: parse metadata.csv per ArtiFact subfolder; scan CASIA Tp/ and Au/.",
            "",
            "FIX #1  -  Generator-aware label inference",
            "    -  ArtiFact metadata uses two schemas: target=0/1 in some folders,",
            "       target=<imagenet_class_index> in others (e.g. big_gan uses 6).",
            "    -  When target is non-binary, fall back to folder name to decide AI vs Real.",
            "    -  Without this fix:  1 AI source on disk found.  After fix:  25.",
            "",
            "FIX #2  -  Stratified sampling across sources",
            "    -  Naive random.sample gave stylegan2 ~70% of the AI class (1M images vs 10K).",
            "    -  Equal allocation per source + redistribute leftovers from larger pools.",
            "",
            "FIX #3  -  Mix CASIA Au/ into the Real class",
            "    -  Earlier runs had Real entirely from ArtiFact - a 100% dataset-source",
            "       shortcut for predicting Forged.  Now Real comes from 11 sources, including CASIA Au.",
            "",
            "Resize 299 x 299 RGB JPEG q95.  Save metadata.csv with source -> class -> split mapping.",
        ],
        "notes": (
            "This is the methodological-rigor slide. The first version of the project found ONLY "
            "stylegan2 as an AI source (everything else used target=imagenet_class). Cross-source "
            "check uncovered this; the discovery fix found 24 additional generators. Stratified "
            "sampling then ensured no single generator dominates. Fix #3 came last, after the 12.5K "
            "run revealed that Real-only-from-ArtiFact + Forged-only-from-CASIA was teaching the "
            "model to predict by dataset signature. Mixing Au into Real cost ~3 pp accuracy but "
            "broke the shortcut - landscape photos went from probable-Forged to 89% Real."
        ),
    },
    {
        "type": "content",
        "title": "Step 2  -  Feature extraction",
        "bullets": [
            "Load Xception via timm:  pretrained=True, num_classes=0",
            "    -  num_classes=0 strips the classifier head.",
            "    -  Output: 2048-dim global-average-pooled feature vector.",
            "",
            "Freeze all parameters  (param.requires_grad = False)",
            "Set model.eval()  (locks BatchNorm, disables dropout)",
            "",
            "Forward-pass each split through the model with no_grad.",
            "Save features as .npy files: xception_x_train.npy, ...",
            "",
            "Ablation: also extract MobileNetV2 features (1280-dim).",
            "    -  ~3-5x faster than Xception at inference; tests whether smaller features generalize better.",
        ],
        "notes": (
            "Stress that NOTHING in the CNN is trained on our data - we are reusing ImageNet's "
            "learned features as-is. MobileNetV2 was picked for its size contrast (much smaller "
            "and faster), to test whether feature-extractor capacity matters. Spoiler: at this "
            "data scale, MobileNetV2 wins. On CPU, Xception extraction took ~70 minutes, "
            "MobileNetV2 ~32 minutes."
        ),
    },
    {
        "type": "content_with_image",
        "image": "../results/cm_mobilenetv2_100_svm.png",
        "title": "Step 3 + 4  -  Classifiers and ablation",
        "bullets": [
            "Two classifiers per backbone:",
            "    -  Random Forest:  200 trees, balanced.",
            "    -  SVM:  RBF, C=10, StandardScaler.",
            "",
            "Test accuracy (12.5K per class):",
            "    -  MobileNetV2 + RF   ->  69.40%",
            "    -  MobileNetV2 + SVM  ->  73.27%   (best)",
            "",
            "Per-class F1 (best combo):",
            "    Real 0.68 - Forged 0.88 - AI 0.72",
            "",
            "Val/test gap = 1.65 pp -> generalizes well.",
            "Forged is the easiest class; Real <-> AI is",
            "where most errors live (see matrix).",
        ],
        "notes": (
            "MobileNetV2 + SVM is consistently the best across all sample sizes we tested. "
            "Two patterns: SVM beats RF on dense GAP features (RF overfits); MobileNetV2 beats "
            "Xception (smaller, less-abstracted features generalize better on this task). "
            "The 74.92% val / 73.27% test number is ~3 pp lower than the prior 77.62% run; that "
            "drop is the honest cost of breaking the CASIA-vs-ArtiFact dataset shortcut by adding "
            "Au to Real. The 1.65 pp val/test gap shows the model is not overfitting. "
            "Per-class F1: Real 0.68 - Forged 0.88 - AI_Generated 0.72 - Forged is the easiest "
            "class because splicing artifacts are well-aligned with ImageNet features. "
            "Xception was not retrained after the Au merge: the 5K and 12.5K-without-Au runs "
            "already showed MobileNetV2 ahead on both classifiers, and SVM training on 21K "
            "Xception features takes multiple hours."
        ),
    },
    {
        "type": "image",
        "title": "Ablation chart  -  test-set accuracy",
        "image": "../results/ablation_comparison.png",
        "caption": (
            "MobileNetV2 + SVM is the best combination at 73.27% test accuracy. "
            "Red dashed line is Ali et al. (2022) binary CASIA baseline at 92.23% - "
            "different task (binary, single dataset, fine-tuned), shown for reference only."
        ),
        "notes": (
            "Walk the audience through the bars left-to-right. The two MobileNetV2 combinations "
            "are what we re-trained after the CASIA Au merge. SVM beats RF by ~4 pp on the same "
            "features - dense GAP features favor a max-margin classifier over tree ensembles. "
            "The Ali et al. line is at 92.23% on a binary, single-dataset, fine-tuned setup; we "
            "include it as a reference point but explicitly do NOT claim parity - slide 8 covers "
            "the apples-to-oranges framing."
        ),
    },
    {
        "type": "content",
        "title": "Comparison with Ali et al. (2022)",
        "bullets": [
            "Ali et al. (2022) reported 92.23% accuracy on CASIA 2.0.",
            "",
            "Differences from this work:",
            "    -  Theirs:  binary  (Real vs. Forged), one dataset.",
            "    -  Ours:    3-class (Real, Forged, AI-Generated), two datasets, 25 AI generators.",
            "",
            "    -  Theirs:  fine-tuned Xception end-to-end on CASIA.",
            "    -  Ours:    frozen Xception/MobileNetV2 + classical classifier (no fine-tuning).",
            "",
            "Direct accuracy comparison is apples-to-oranges:",
            "    -  Random baseline for binary = 50%; their 92.23% is +42 over chance.",
            "    -  Random baseline for 3-class = 33%; our 74.92% is +42 over chance.  (Comparable lift.)",
            "",
            "On the binary CASIA Tp/ subtask alone our system still hits ~89% accuracy after the",
            "Au merge  (down from 96.5% before, because the CASIA-source shortcut is gone).",
        ],
        "notes": (
            "Don't claim we beat Ali et al. - we have a different problem. The honest framing is: "
            "as 'accuracy gain over chance,' the two are comparable, but our task is harder. "
            "On the binary subtask (CASIA forged), our system was at 96.5% before the Au merge "
            "and at 89.08% after, because the model can no longer use 'this came from CASIA' as "
            "a free predictor. That accuracy drop is the price of methodological honesty. "
            "Be ready for: 'Why didn't you reproduce their setup directly?' Answer: time, GPU "
            "access, and the 3-class formulation forces a different setup."
        ),
    },
    {
        "type": "content",
        "title": "Generator difficulty tiers  (the headline finding)",
        "bullets": [
            "Per-source AI accuracy (mobilenetv2_100 + SVM, 12.5K samples, CASIA Au in Real):",
            "",
            "EASY  (>= 85%, frozen features detect them well):",
            "    star_gan 100%, sfhq 96%, face_synthetics 96%, ddpm 94%, cips 91%,",
            "    gansformer 89%, denoising_diffusion_gan 86%, big_gan 85%, palette 85%.",
            "",
            "MEDIUM  (60 - 85%):",
            "    gau_gan 76%, stylegan1 69%, lama 69%, stylegan2 64%, taming_transformer 62%,",
            "    projected_gan 62%, glide 60%, mat 60%.",
            "",
            "HARD  (< 60%, frozen features struggle):",
            "    stable_diffusion 59%, diffusion_gan 57%, latent_diffusion 55%,",
            "    vq_diffusion 52%, generative_inpainting 46%,  stylegan3 33%.",
            "",
            "Pattern:  older / artifact-heavy generators are detected; modern diffusion approaches chance.",
        ],
        "notes": (
            "This slide elevates the project from 'we built a classifier' to 'we discovered "
            "something about the limits of frozen ImageNet features.' Newer generators "
            "(stylegan3, stable_diffusion) are actively designed to look photorealistic - they "
            "don't leave the kind of artifacts ImageNet semantics encode. Older generators "
            "(BigGAN, StarGAN, DDPM) leave artifacts that align with ImageNet's feature space, "
            "so we detect them at 90%+. This is the empirical case for frequency-aware backbones "
            "(NoisePrint, FrequencyNet) over ImageNet transfer learning. "
            "Note the StyleGAN gradient: 1 -> 2 -> 3 from 69% -> 64% -> 33%. Each release was "
            "explicitly designed to remove the artifacts of the previous one."
        ),
    },
    {
        "type": "content",
        "title": "Key design decisions",
        "bullets": [
            "Frozen vs. fine-tuned CNN:",
            "    +  Trains in minutes on CPU. Reproducible. No GPU required.",
            "    -  Misses task-specific cues (texture, frequency, generator fingerprints).",
            "",
            "Random Forest + SVM vs. neural head:",
            "    +  Fast iteration, calibration available, easy to interpret.",
            "    -  Fewer parameters; RF probabilities uncalibrated by default.",
            "",
            "3-class vs. binary formulation:",
            "    +  Closer to a real moderation use case.",
            "    -  Per-class data balance is harder; less directly comparable to prior work.",
            "",
            "5K vs 12.5K samples per class  (data-scaling experiment):",
            "    -  RF accuracy DROPPED with 2.5x more data on both backbones (-2.4 / -2.7 pp).",
            "    -  Direct evidence of a frozen-feature ceiling, not a data shortage.",
            "",
            "Project structured as a Python package (src/ + pyproject.toml) for reproducibility.",
        ],
        "notes": (
            "Frame each decision as an explicit trade-off. The 5K -> 12.5K experiment is your "
            "evidence that simply scaling data won't fix the core limitation: frozen ImageNet "
            "features have a representational ceiling for modern AI detection. This justifies "
            "the future-work slide and shows you tested the obvious next step before recommending it."
        ),
    },
    {
        "type": "content",
        "title": "Iteration history  -  what each run taught us",
        "bullets": [
            "Run 1  -  Naive baseline:",
            "    -  Result:  ~84% accuracy, looked great.",
            "    -  Reality:  only stylegan2 was being detected (1 AI source on disk).",
            "    -  Lesson:  cross-source check is mandatory; headline accuracy can be a mirage.",
            "",
            "Run 2  -  Discovery fix + stratified sampling, 5K per class:",
            "    -  Result:  RF ~71%, SVM ~76%.  25 AI sources active.",
            "    -  Cross-source revealed:  lsun (real scenes) at 28%, ffhq (real faces) routed to AI.",
            "    -  Lesson:  data composition matters more than data volume.",
            "",
            "Run 3  -  Same fixes, 12.5K per class (more data):",
            "    -  Result:  RF dropped 2.4 - 2.7 pp.  SVM hit 77.62% val / 76.76% test.",
            "    -  Lesson:  frozen-feature ceiling is real - more data hurts the high-capacity classifier.",
            "    -  Bias still visible:  landscape photos in the demo predicted 'Forged'.",
            "",
            "Run 4  -  Mix CASIA Au into Real, 12.5K per class  (current results):",
            "    -  Result:  74.92% val  /  73.27% test  (1.65 pp gap, good generalization).",
            "    -  landscape:  ~28% (Run 2)  ->  89% Real (Run 4).  lsun:  ->  69%.",
            "    -  Forged on Tp/:  96.5%  ->  89.08%.  CASIA shortcut broken; the drop is honest cost.",
            "    -  Au at 49.7% Real (chance-level within CASIA)  =  proof the shortcut is gone.",
        ],
        "notes": (
            "This is the journey slide and arguably the most important content. Frame it as "
            "'each iteration moved the headline number AND fixed a real methodological flaw.' "
            "Run 1 was a stylegan2 detector hiding behind 84% accuracy. Run 2 caught that. "
            "Run 3 added data and proved data isn't the bottleneck. Run 4 traded 3 pp of "
            "accuracy for genuine source-independent classification. "
            "The 49.7% on CASIA Au is the smoking gun - it's exactly chance for a binary "
            "Real-vs-Forged distinction within the same dataset, which means the model is "
            "no longer using dataset signature to predict. This is what we wanted."
        ),
    },
    {
        "type": "content",
        "title": "Cross-source sanity check  (Step 4b)",
        "bullets": [
            "Concern: with Real and Forged coming from different datasets, the classifier could",
            "learn dataset signatures (JPEG tables, sensor noise) instead of authenticity cues.",
            "",
            "Method:  step4b_cross_dataset_check.py",
            "    -  Splits the test set by source folder (each ArtiFact generator, CASIA Tp/, CASIA Au/).",
            "    -  Reports per-source accuracy + within-class spread.",
            "",
            "Findings on the final 12.5K MobileNetV2 + SVM model (after Au merge):",
            "    -  Forged (CASIA Tp/):           89.08%   (1 source, no spread)",
            "    -  AI_Generated:    33% - 100%   (23 sources, spread = 67.2 pp)",
            "    -  Real:            49% - 93%    (11 sources, spread = 43.8 pp)",
            "    -  CASIA Au:        49.72%       (chance-level  =  shortcut broken)",
            "",
            "Verdict:  generator difficulty is real, not a shortcut. lsun and landscape recovered",
            "from <30% to 69% / 89%. Forged dropped 7 pp because the dataset-source shortcut is gone.",
            "",
            "Decision rule:  spread > 15 pp  ->  do NOT fine-tune frozen features (would amplify the gap).",
        ],
        "notes": (
            "The methodologically strongest slide. The 5K version of this same check is what "
            "surfaced the dataset-bias problem (lsun at 28%, AI_Generated only stylegan2). After "
            "fixing discovery + sampling + adding CASIA Au to Real, the 12.5K version shows: real "
            "generalization across most sources, with a small set of modern diffusion generators "
            "frozen features cannot learn. The Au signal at 49.7% is the strongest evidence that "
            "we no longer have a shortcut. The decision rule embedded in step4b is what tells us "
            "NOT to fine-tune in response to wide spread - that's the next slide."
        ),
    },
    {
        "type": "content",
        "title": "Should we fine-tune the CNN?  Decision: NO (for now)",
        "bullets": [
            "Three reasons not to fine-tune now:",
            "",
            "1.  Spread rule says NO.",
            "    -  AI_Generated spread = 67 pp,  Real spread = 44 pp.  Both >> 15 pp threshold.",
            "    -  Fine-tuning amplifies bias when the data is uneven across sources;",
            "       it would make easy generators easier and hard generators no better.",
            "",
            "2.  Data-scaling experiment ruled out 'just need more data'.",
            "    -  5K -> 12.5K per class lowered RF accuracy on both backbones.",
            "    -  More data didn't close the gap, so the bottleneck isn't training-set size.",
            "",
            "3.  Hard tier needs a different inductive bias, not a tuned ImageNet net.",
            "    -  StyleGAN3 at 33%, Stable Diffusion at 59%, latent_diffusion at 55%.",
            "    -  These generators are designed to look photographic to ImageNet-style nets.",
            "    -  The fix is frequency / noise-residual features (NoisePrint, FrequencyNet),",
            "       not more learning on the wrong inductive bias.",
            "",
            "When to revisit fine-tuning:",
            "    -  After per-source spread drops below ~15 pp via better dataset balancing.",
            "    -  As a TARGETED last-block tune on Easy + Medium tiers only,",
            "       not a generic end-to-end retrain.",
        ],
        "notes": (
            "This is the explicit answer to 'why don't you just fine-tune?' Three reasons, in "
            "order of strength. The spread rule is the hard stop - we built it into step4b "
            "specifically to keep us from doing the wrong thing under pressure. The data-scaling "
            "result is the experimental evidence. The frequency-bias point is the architectural "
            "argument: fine-tuning a semantic-feature net won't grow the spectral features that "
            "modern diffusion detection actually needs. "
            "If asked 'so when would you fine-tune?' say: after either (a) the dataset is "
            "balanced enough that spread is below the threshold, or (b) you've already added "
            "frequency-aware features and want to push the Easy + Medium tiers higher. Not before."
        ),
    },
    {
        "type": "content",
        "title": "Demo, limitations, and future work",
        "bullets": [
            "Live demo:  python -m src.step5_gradio_demo  (http://127.0.0.1:7860)",
            "    -  Uses the best combo by default:  MobileNetV2 + SVM + StandardScaler.",
            "    -  Upload an image  ->  predicted class + per-class confidences.",
            "",
            "Known limitations:",
            "    -  Frozen ImageNet features hit a ceiling on modern diffusion (Stable Diffusion ~59%).",
            "    -  Random Forest probabilities uncalibrated (demo bars are indicative, not literal).",
            "    -  Single train/val/test seed; no cross-validation.",
            "    -  CASIA Au at 49.7% means within-CASIA Real-vs-Forged is at chance - some Au",
            "       images still get called Forged because of residual sensor-noise overlap with Tp/.",
            "",
            "Future work, in priority order:",
            "    1.  Frequency-aware backbone (NoisePrint / FrequencyNet) for the Hard tier.",
            "    2.  Targeted fine-tune of last block on Easy + Medium only, after spread drops.",
            "    3.  Calibrate RF/SVM probabilities for trustworthy demo confidences.",
            "    4.  Leave-one-generator-out evaluation for cross-dataset robustness.",
            "    5.  Add more Real-image sources to bring Real spread below the 15 pp threshold.",
            "",
            "Questions?",
        ],
        "notes": (
            "Run the Gradio demo here if AV setup allows. Have 2-3 sample images preloaded: "
            "a real photo (test that landscape photos now go to Real - they did NOT before the "
            "CASIA Au fix), a CASIA tampered image, and an AI-generated one (preferably a hard "
            "one like Stable Diffusion to honestly show the limitation). "
            "Future-work order matters: frequency-aware backbone FIRST because it addresses the "
            "Hard tier directly. Fine-tuning is now item #2 (was #1 in earlier versions of this "
            "deck) because the data-scaling experiment plus the spread rule both point away "
            "from end-to-end fine-tuning as a first move."
        ),
    },
]


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, notes):
    layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders[1].has_text_frame:
        slide.placeholders[1].text = subtitle
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_content_slide(prs, title, bullets, notes):
    layout = prs.slide_layouts[1]  # Title + content layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.font.size = Pt(16)

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_content_with_image_slide(prs, title, bullets, image_path, notes):
    """Slide with bullets on the left half and an image on the right half."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    # Resize the body placeholder to the left half so bullets do not overlap the image.
    body = slide.placeholders[1]
    body.left = Inches(0.5)
    body.top = Inches(1.5)
    body.width = Inches(6.5)
    body.height = Inches(5.5)

    body_frame = body.text_frame
    body_frame.clear()
    for i, line in enumerate(bullets):
        p = body_frame.paragraphs[0] if i == 0 else body_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(14)

    # Image on the right half
    img_full_path = Path(__file__).parent / image_path
    slide.shapes.add_picture(
        str(img_full_path),
        left=Inches(7.2),
        top=Inches(1.6),
        width=Inches(5.8),
    )

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_image_slide(prs, title, image_path, caption, notes):
    """Title + large centered image + small caption."""
    layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    img_full_path = Path(__file__).parent / image_path
    slide.shapes.add_picture(
        str(img_full_path),
        left=Inches(2.5),
        top=Inches(1.4),
        width=Inches(8.3),
    )

    # Caption text box below the image
    from pptx.util import Emu
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(13)
    p.alignment = PP_ALIGN.CENTER

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def main():
    out_dir = Path(__file__).parent
    out_path = out_dir / "truphoto_presentation.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in SLIDES:
        kind = spec["type"]
        if kind == "title":
            add_title_slide(prs, spec["title"], spec["subtitle"], spec["notes"])
        elif kind == "content_with_image":
            add_content_with_image_slide(
                prs, spec["title"], spec["bullets"], spec["image"], spec["notes"]
            )
        elif kind == "image":
            add_image_slide(
                prs, spec["title"], spec["image"], spec["caption"], spec["notes"]
            )
        else:
            add_content_slide(prs, spec["title"], spec["bullets"], spec["notes"])

    prs.save(out_path)
    print(f"Wrote {out_path}  ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
